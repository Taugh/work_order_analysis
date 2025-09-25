# Scripts/slide_updater.py
# ---------------------------------------------------------------
# This module updates PowerPoint charts for work order analysis.
# It requires a PowerPoint template, pandas DataFrames/dicts with summary/group data,
# and chart names that match those in the template.
# Outputs: Updates charts in the PowerPoint file and saves the result to the outputs/presentations directory.
# ---------------------------------------------------------------

from pptx import Presentation
from pptx.chart.data import CategoryChartData
import pandas as pd
import logging
import os

logging.basicConfig(level=logging.INFO)

def update_pm_missed_chart(prs, slide_index, chart_name, chart_data):
    """
    Updates the PM Missed Chart on a specified slide in the PowerPoint presentation.

    Requirements:
        - prs: Presentation object loaded from a template.
        - slide_index: Index of the slide containing the chart.
        - chart_name: Name of the chart shape to update (must match template).
        - chart_data: Dict with keys 'months', 'due', 'complete', 'missed'.

    Output:
        - Replaces chart data for the specified chart on the slide.
        - No return value; modifies prs in place.
    """
    slide = prs.slides[slide_index]
    chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, "chart") and shape.name == chart_name:
            chart_shape = shape
            break
    if not chart_shape:
        logging.warning(f"Chart '{chart_name}' not found on slide {slide_index+1}.")
        return
    cat_data = CategoryChartData()
    cat_data.categories = chart_data["months"]
    cat_data.add_series("Due", chart_data["due"])
    cat_data.add_series("Completed", chart_data["complete"])
    cat_data.add_series("Missed", chart_data["missed"])
    chart_shape.chart.replace_data(cat_data)

def update_group_charts(prs, group_data, slide_index):
    """
    Updates group-level missed work order charts on a specified slide.

    Requirements:
        - prs: Presentation object loaded from a template.
        - group_data: Dict with keys 'groups', 'missed', 'missed_percent'.
        - slide_index: Index of the slide containing the charts.
        - Chart shapes named "Qty Missed by Group" and "% Missed by Group" must exist.

    Output:
        - Replaces chart data for both charts on the slide.
        - No return value; modifies prs in place.
    """
    slide = prs.slides[slide_index]
    qty_chart_shape = None
    percent_chart_shape = None
    for shape in slide.shapes:
        if hasattr(shape, "chart"):
            if shape.name == "Qty Missed by Group":
                qty_chart_shape = shape
            elif shape.name == "% Missed by Group":
                percent_chart_shape = shape
    filtered = [
        (g, m, p)
        for g, m, p in zip(group_data["groups"], group_data["missed"], group_data["missed_percent"])
        if m > 0
    ]
    if not filtered:
        logging.warning("No groups with missed work orders.")
        return
    groups, missed, missed_percent = zip(*filtered)
    if qty_chart_shape:
        qty_data = CategoryChartData()
        qty_data.categories = groups
        qty_data.add_series("Missed", missed)
        qty_chart_shape.chart.replace_data(qty_data)
    else:
        logging.warning("Qty Missed by Group chart not found.")
    if percent_chart_shape:
        percent_data = CategoryChartData()
        percent_data.categories = groups
        percent_data.add_series("% Missed", missed_percent)
        percent_chart_shape.chart.replace_data(percent_data)
    else:
        logging.warning("% Missed by Group chart not found.")

# Example usage:
# ---------------------------------------------------------------
# Loads a template presentation, updates charts with summary/group data,
# and saves the updated presentation to outputs/presentations.
# ---------------------------------------------------------------

summary_data = {
    "months": ["Jan", "Feb", "Mar", "Apr"],
    "due": [50, 55, 60, 58],
    "complete": [45, 52, 55, 53],
    "missed": [5, 3, 5, 5]
}
group_data = {
    "groups": ["Ops", "Tech", "Admin"],
    "missed": [6, 8, 3],
    "missed_percent": [12.5, 18.0, 7.1],
    "reporting_month": "2025-06"
}

prs = Presentation("template.pptx")
update_pm_missed_chart(prs, slide_index=1, chart_name="PM Missed Chart", chart_data=summary_data)
update_group_charts(prs, group_data, slide_index=2)

for i, shape in enumerate(prs.slides[1].shapes):
    print(f"Slide 2 Shape {i}: type={shape.shape_type}, name={getattr(shape, 'name', None)}")
for i, shape in enumerate(prs.slides[2].shapes):
    print(f"Slide 3 Shape {i}: type={shape.shape_type}, name={getattr(shape, 'name', None)}")

# Save output presentation
last_month_label = pd.to_datetime(group_data["reporting_month"], format="%b-%y").strftime("%b-%Y")
output_ppt_path = f"outputs/presentations/group_slide_deck_{last_month_label}.pptx"
os.makedirs(os.path.dirname(output_ppt_path), exist_ok=True)
prs.save(output_ppt_path)
logging.info(f"✅ Saved updated charts to {output_ppt_path}")