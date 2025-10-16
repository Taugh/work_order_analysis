# ---------------------------------------------------------------
# slide_generator.py
#
# Purpose:
#   Generates and updates PowerPoint slides for work order governance reporting.
#
# Requirements:
#   - Input: pandas DataFrames containing summary, group, and monthly work order data.
#   - PowerPoint template: 'data/templates/governance_slide_template.pptx' with named charts/shapes.
#   - DataFrames must include columns such as 'Month', 'Due', 'Completed', 'Missed', 'group', etc.
#
# Output:
#   - Creates and updates slides with summary metrics, charts, and tables.
#   - Saves the final presentation to 'outputs/presentations/governance_slide.pptx' (or specified path).
#
# Notes:
#   - Functions are provided for both creating new slides and updating existing charts/tables.
#   - All chart and table updates require matching shape/chart names in the PowerPoint template.
#   - Designed for use in both CLI and GUI workflows.
# ---------------------------------------------------------------

# scripts/slide_generator.py

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import pandas as pd
import numpy as np
from datetime import datetime
import traceback

def create_full_governance_deck(by_month_df, late_df, disposition_df, by_group_df, filename=None):
    """
    Creates a complete governance presentation with all charts and data
    
    Parameters:
    - by_month_df: DataFrame with monthly trend data (from trend_df in main.py)
    - late_df: DataFrame with late work orders
    - disposition_df: DataFrame with disposition data (can be None)
    - by_group_df: DataFrame with group performance data
    - filename: Optional custom filename
    """
    try:
        # FIX: Correct template filename
        template_path = "data/templates/governance_slide_template.pptx"
        prs = Presentation(template_path)
        print(f"✅ Loaded template: {template_path}")
        
        # Generate filename if not provided
        if filename is None:
            current_month = datetime.now().strftime("%b")  # e.g., "Sep"
            filename = f"governance_slide_{current_month}.pptx"
            print(f"🔍 DEBUG: Auto-generated PowerPoint filename: {filename}")
        
        # FIX: Add slide 0 processing for monthly and YTD summaries
        update_summary_slide(prs, by_month_df, late_df, slide_index=0)
        
        # Update all slides with data
        update_missed_by_month_chart(prs, by_month_df, slide_index=1)
        update_missed_disposition_chart(prs, disposition_df, slide_index=2)
        update_group_charts(prs, by_group_df, slide_index=3)
        
        # Save the presentation
        output_path = f"outputs/presentations/{filename}"
        prs.save(output_path)
        print(f"✅ Full governance deck saved to: {output_path}")
        
        return output_path
        
    except Exception as e:
        print(f"❌ Error creating governance deck: {e}")
        traceback.print_exc()
        return None

def update_summary_slide(prs, by_month_df, late_df, slide_index=0):
    """
    Updates slide 0 with monthly summary and YTD summary
    """
    try:
        slide = prs.slides[slide_index]
        
        # Calculate current month and YTD data
        today = datetime.now()
        current_month = today.strftime("%b-%y")
        current_year = today.year % 100  # Get 2-digit year
        
        # Monthly Summary (Previous Month)
        prev_month = (today.replace(day=1) - pd.DateOffset(months=1)).strftime("%b-%y")
        monthly_data = by_month_df[by_month_df['report_month'] == prev_month]
        
        if not monthly_data.empty:
            monthly_generated = monthly_data['generated'].iloc[0]
            monthly_completed = monthly_data['completed'].iloc[0]
            monthly_missed = monthly_data['missed'].iloc[0]
            monthly_completion_rate = (monthly_completed / monthly_generated * 100) if monthly_generated > 0 else 0
        else:
            monthly_generated = monthly_completed = monthly_missed = 0
            monthly_completion_rate = 0
        
        # YTD Summary (Current Year)
        ytd_data = by_month_df[by_month_df['report_month'].str.endswith(f'-{current_year:02d}')]
        ytd_generated = ytd_data['generated'].sum()
        ytd_completed = ytd_data['completed'].sum()
        ytd_missed = ytd_data['missed'].sum()
        ytd_completion_rate = (ytd_completed / ytd_generated * 100) if ytd_generated > 0 else 0
        
        # Count late work orders for summary
        late_count = len(late_df) if late_df is not None and not late_df.empty else 0
        
        # FIX: Ensure the title is preserved
        # First, find and preserve the title
        title_shape = None
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                text_content = shape.text_frame.text
                # Check if this is likely the title (position and content)
                if (hasattr(shape, 'top') and shape.top < Inches(1.5) and 
                    ("PM Monthly" in text_content or "Summary" in text_content or 
                     "YTD" in text_content and len(text_content) < 100)):
                    title_shape = shape
                    # Set the correct title
                    title_shape.text_frame.clear()
                    p = title_shape.text_frame.paragraphs[0]
                    p.text = "PM Monthly and YTD and Summaries Completion Rates"
                    p.font.name = 'Aptos'
                    p.font.size = Pt(32)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 76, 153)  # Blue color
                    print("✅ Updated slide title")
                    break
        
        # Find and update content text boxes (excluding title)
        monthly_updated = False
        ytd_updated = False
        
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame and shape != title_shape:
                text_content = shape.text_frame.text
                
                # Monthly Summary Text Box (not in title area)
                if (("Monthly Summary" in text_content or "Previous Month" in text_content or 
                     "Month Summary" in text_content) and 
                    hasattr(shape, 'top') and shape.top > Inches(1.5) and not monthly_updated):
                    new_text = f"Monthly Summary ({prev_month}):\nGenerated: {monthly_generated:,}\nCompleted: {monthly_completed:,}\nMissed: {monthly_missed:,}\nCompletion Rate: {monthly_completion_rate:.1f}%"
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = new_text
                    p.font.name = 'Aptos'
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 102, 51)  # Greenish color
                    monthly_updated = True
                    print(f"✅ Updated Monthly Summary for {prev_month}")
                
                # YTD Summary Text Box (not in title area)
                elif (("YTD Summary" in text_content or "Year to Date" in text_content) and 
                      hasattr(shape, 'top') and shape.top > Inches(1.5) and not ytd_updated):
                    new_text = f"YTD Summary (20{current_year}):\nGenerated: {ytd_generated:,}\nCompleted: {ytd_completed:,}\nMissed: {ytd_missed:,}\nCompletion Rate: {ytd_completion_rate:.1f}%"
                    shape.text_frame.clear()
                    p = shape.text_frame.paragraphs[0]
                    p.text = new_text
                    p.font.name = 'Aptos'
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 102, 51)  # Greenish color
                    ytd_updated = True
                    print(f"✅ Updated YTD Summary for 20{current_year}")
        
        # If no existing content text boxes found, create them
        if not monthly_updated:
            # Add Monthly Summary text box
            left = Inches(0.5)
            top = Inches(2)  # Below title
            width = Inches(5)
            height = Inches(3.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = f"Monthly Summary ({prev_month}):\nGenerated: {monthly_generated:,}\nCompleted: {monthly_completed:,}\nMissed: {monthly_missed:,}\nCompletion Rate: {monthly_completion_rate:.1f}%"
            text_frame.paragraphs[0].font.size = Pt(16)
            text_frame.paragraphs[0].font.bold = True
            print(f"📝 Created new Monthly Summary text box for {prev_month}")
        
        if not ytd_updated:
            # Add YTD Summary text box
            left = Inches(5.5)
            top = Inches(2)  # Below title
            width = Inches(5)
            height = Inches(3.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = f"YTD Summary (20{current_year}):\nGenerated: {ytd_generated:,}\nCompleted: {ytd_completed:,}\nMissed: {ytd_missed:,}\nCompletion Rate: {ytd_completion_rate:.1f}%\nLate Orders: {late_count:,}"
            text_frame.paragraphs[0].font.size = Pt(16)
            text_frame.paragraphs[0].font.bold = True
            print(f"📝 Created new YTD Summary text box for 20{current_year}")
        
        print("✅ Summary slide (slide 0) updated successfully")
        
    except Exception as e:
        print(f"❌ Error updating summary slide: {e}")
        traceback.print_exc()

def update_missed_by_month_chart(prs, by_month_df, slide_index=1):
    """
    Updates the missed by month chart on slide 1 and adds stoplight table
    """
    slide = prs.slides[slide_index]
    
    # Debug: Print DataFrame info
    print("🔍 DEBUG: by_month_df columns:", by_month_df.columns.tolist())
    print("🔍 DEBUG: by_month_df shape:", by_month_df.shape)
    print("🔍 DEBUG: by_month_df head:")
    print(by_month_df.head())
    
    # Find and update the chart
    chart_found = False
    text_box_found = False
    
    # Handle different possible column names for months
    month_column = None
    if 'report_month' in by_month_df.columns:
        month_column = 'report_month'
    elif 'Month' in by_month_df.columns:
        month_column = 'Month'
    elif by_month_df.index.name == 'report_month':
        # Month data is in the index
        by_month_df = by_month_df.reset_index()
        month_column = 'report_month'
    else:
        # Try to use the first column or index
        if len(by_month_df.columns) > 0:
            month_column = by_month_df.columns[0]
        else:
            print("❌ Could not find month column in DataFrame")
            return
    
    # Map column names to standardized names
    column_mapping = {
        'generated': ['Due', 'generated', 'Generated'],
        'completed': ['Completed', 'completed'], 
        'missed': ['Missed', 'missed']
    }
    
    # Find the actual column names
    actual_columns = {}
    for standard_name, possible_names in column_mapping.items():
        for possible_name in possible_names:
            if possible_name in by_month_df.columns:
                actual_columns[standard_name] = possible_name
                break
        if standard_name not in actual_columns:
            print(f"❌ Could not find column for '{standard_name}' in DataFrame")
            return
    
    print(f"📅 Using month column: '{month_column}'")
    print(f"📊 Column mapping: {actual_columns}")
    
    for shape in slide.shapes:
        # Update chart with new data
        if hasattr(shape, "chart"):
            chart = shape.chart
            chart_data = CategoryChartData()
            
            # Use identified month column as categories
            categories = by_month_df[month_column].astype(str).tolist()
            chart_data.categories = categories
            
            # Add data series using mapped column names
            # FIX: Correct order - Generated, Missed, Completed
            chart_data.add_series('Missed', by_month_df[actual_columns['missed']].tolist())
            chart_data.add_series('Completed', by_month_df[actual_columns['completed']].tolist())
            chart_data.add_series('Generated', by_month_df[actual_columns['generated']].tolist())
            
            # Replace chart data
            chart.replace_data(chart_data)
            chart_found = True
            print("✅ Missed by Month chart updated")
            break
    
    # Update rolling totals text box
    total_due = by_month_df[actual_columns['generated']].sum()
    total_completed = by_month_df[actual_columns['completed']].sum()
    total_missed = by_month_df[actual_columns['missed']].sum()
    
    print(f"📊 Rolling totals calculated:")
    print(f"  Total Due: {total_due}")
    print(f"  Total Completed: {total_completed}")
    print(f"  Total Missed: {total_missed}")
    
    # Look for existing text box with "Rolling 12-Month" text
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            text_content = shape.text_frame.text
            if "Rolling 12-Month" in text_content or "12-Month" in text_content:
                print("🎯 Found existing Rolling 12-Month text box - updating it")
                print(f"📍 Template text box position: left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}")
                
                # Update the text content
                new_text = f"12-Month Totals:\nDue: {total_due:,}\nCompleted: {total_completed:,}\nMissed: {total_missed:,}"
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = new_text
                p.font.size = Pt(18)
                p.font.bold = True
                
                text_box_found = True
                print(f"✅ Rolling 12-month totals updated in existing text box")
                print(f"✅ Final text frame content: '{shape.text_frame.text}'")
                break
    
    # Add stoplight table (also needs to be updated)
    add_stoplight_table_two_tables(prs, slide_index, by_month_df, month_column, actual_columns)
    
    if not text_box_found:
        print("❌ Could not find existing 'Rolling 12-Month' text box in template")
        print("💡 Make sure your template has a text box containing 'Rolling 12-Month' text")

def add_stoplight_table_two_tables(prs, slide_index=1, by_month_df=None, month_column='Month', actual_columns=None):
    """
    Creates two separate tables:
    1. Header table (1 row x 1 column) for the key
    2. Data table (2 rows x 12 columns) for months and stoplights
    """
    try:
        slide = prs.slides[slide_index]
        
        # Extract months using the identified month column
        if by_month_df is not None and not by_month_df.empty:
            months = by_month_df[month_column].astype(str).tolist()[:12]
            print(f"📅 Using months from chart data: {months}")
        else:
            months = ['Oct-24', 'Nov-24', 'Dec-24', 'Jan-25', 'Feb-25', 'Mar-25',
                      'Apr-25', 'May-25', 'Jun-25', 'Jul-25', 'Aug-25', 'Sep-25']
            print(f"⚠️ Using default months: {months}")
        
        # Ensure exactly 12 months
        while len(months) < 12:
            months.append("")
        months = months[:12]
        
        # === TABLE 1: HEADER TABLE ===
        print("📋 Creating header table...")
        header_left = Inches(0.42)
        header_top = Inches(0.75)
        header_width = Inches(10.17)
        header_height = Inches(0.25)
        
        # Create header table (1 row x 1 column)
        header_table_shape = slide.shapes.add_table(1, 1, header_left, header_top, header_width, header_height)
        header_table = header_table_shape.table
        
        # Style header table
        header_cell = header_table.cell(0, 0)
        header_cell.text = "Stop Lights -  🔴: Missed > 6%,    🟡: MISSED >3% <=6%,    🟢 <=3%"
        header_cell.text_frame.paragraphs[0].font.size = Pt(16)
        header_cell.text_frame.paragraphs[0].font.bold = True
        header_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        header_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
        header_cell.fill.solid()
        header_cell.fill.fore_color.rgb = RGBColor(51, 153, 255)  # Light blue background
        
        print("✅ Header table created")
        
        # === TABLE 2: DATA TABLE ===
        print("📊 Creating data table...")
        data_left = Inches(0.42)
        data_top = Inches(1.1)  # Just below header table
        data_width = Inches(10.17)
        data_height = Inches(0.77)
        
        # Create data table (2 rows x 12 columns)
        data_table_shape = slide.shapes.add_table(2, 12, data_left, data_top, data_width, data_height)
        data_table = data_table_shape.table
        
        # Row 1: Month names
        print("📅 Adding month names...")
        for col, month in enumerate(months):
            cell = data_table.cell(0, col)
            cell.text = str(month) if month else ""
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            # Light blue background for month row
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(51, 153, 255) 
            
        # Row 2: Stoplight indicators
        print("🚦 Adding stoplight indicators...")
        for col, month in enumerate(months):
            cell = data_table.cell(1, col)
            
            if month:  # Only calculate if we have a valid month
                missed_percentage = calculate_performance_metric(by_month_df, month, month_column, actual_columns)
                
                # Determine stoplight color based on missed percentage
                if missed_percentage > 6:
                    stoplight = "🔴"  # RED: Missed > 6%
                    bg_color = RGBColor(102, 0, 0)  # Light red
                elif missed_percentage > 3:
                    stoplight = "🟡"  # YELLOW: Missed >3% <=6%
                    bg_color = RGBColor(204, 204, 0)  # Light yellow
                else:
                    stoplight = "🟢"  # GREEN: <=3%
                    bg_color = RGBColor(0, 204, 0)  # Light green
            else:
                # Empty cell for padding
                stoplight = ""
                bg_color = RGBColor(255, 255, 255)  # White background
            
            cell.text = stoplight
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
        
        # Style both tables
        header_table.first_row = False
        header_table.first_col = False
        data_table.first_row = False
        data_table.first_col = False
        
        print("✅ Two-table stoplight display created successfully")
        return data_table
     
    except Exception as e:
        print(f"❌ Error creating two-table stoplight display: {e}")
        traceback.print_exc()
        return None

def calculate_performance_metric(by_month_df, month, month_column='Month', actual_columns=None):
    """
    Calculate missed percentage for stoplight determination
    Formula: Missed / (Missed + Completed) * 100
    Returns missed percentage (0-100)
    """
    if by_month_df is None or actual_columns is None:
        return 2  # Default value for testing (GREEN)
    
    # Find the row for this month
    month_data = by_month_df[by_month_df[month_column] == month]
    if not month_data.empty:
        missed = month_data[actual_columns['missed']].iloc[0]
        completed = month_data[actual_columns['completed']].iloc[0]
        
        # Calculate missed percentage: Missed / (Missed + Completed)
        total_wo_processed = missed + completed
        missed_percentage = (missed / total_wo_processed) * 100 if total_wo_processed > 0 else 0
        
        print(f"📊 {month}: Missed={missed}, Completed={completed}, Missed%={missed_percentage:.2f}%")
        return missed_percentage
    
    return 0  # Default if month not found

def update_missed_disposition_chart(prs, disposition_df, slide_index=2):
    """
    Updates the missed disposition chart on slide 2
    """
    try:
        if disposition_df is None or disposition_df.empty:
            print("⚠️ No disposition data provided - skipping disposition chart")
            return
        
        slide = prs.slides[slide_index]
        
        print(f"🔍 DEBUG: disposition_df columns: {disposition_df.columns.tolist()}")
        print(f"🔍 DEBUG: disposition_df shape: {disposition_df.shape}")
        
        # FIX: Sort months chronologically instead of alphabetically
        def month_sort_key(month_str):
            """Convert month string to sortable format"""
            try:
                # Parse "MMM-YY" format (e.g., "Oct-24")
                month_date = pd.to_datetime(month_str, format='%b-%y')
                return month_date
            except:
                return pd.to_datetime('1900-01-01')  # Default for invalid dates
        
        # Sort disposition_df by month chronologically
        disposition_df_sorted = disposition_df.copy()
        disposition_df_sorted['sort_key'] = disposition_df_sorted['report_month'].apply(month_sort_key)
        disposition_df_sorted = disposition_df_sorted.sort_values('sort_key').drop('sort_key', axis=1)
        
        print(f"🔍 DEBUG: Months after sorting: {disposition_df_sorted['report_month'].tolist()}")
        
        # Map column names
        disposition_columns = {
            'month': 'report_month',
            'closed': 'closed',
            'awaiting_qa': 'awaiting_qa',
            'awaiting_dept': 'awaiting_dept'
        }
        
        print(f"📊 Disposition column mapping: {disposition_columns}")
        
        # Find existing chart and remove it first
        for shape in slide.shapes:
            if hasattr(shape, "chart"):
                chart_position = {
                    'left': shape.left,
                    'top': shape.top, 
                    'width': shape.width,
                    'height': shape.height
                }
                print(f"🎯 Found chart at position: {chart_position}")
                
                # Remove the existing chart
                slide.shapes._spTree.remove(shape._element)
                print("🗑️ Removed problematic chart")
                break
        
        # Create new chart
        try:
            # Chart position and size
            chart_left = Inches(1)
            chart_top = Inches(1.25)
            chart_width = Inches(10)
            chart_height = Inches(5)
            
            # Prepare chart data
            chart_data = CategoryChartData()
            chart_data.categories = disposition_df_sorted[disposition_columns['month']].astype(str).tolist()
            chart_data.add_series('Closed', disposition_df_sorted[disposition_columns['closed']].tolist())
            chart_data.add_series('Awaiting QA', disposition_df_sorted[disposition_columns['awaiting_qa']].tolist())
            chart_data.add_series('Awaiting Dept', disposition_df_sorted[disposition_columns['awaiting_dept']].tolist())
            
            # Add chart
            chart_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_STACKED, 
                chart_left, chart_top, chart_width, chart_height, chart_data
            )
            chart = chart_shape.chart
            
            # Apply styling
            chart.chart_style = 8
            print("✅ Applied Style 8 to chart")
            
            # Color the series
            series_colors = {
                'Closed': '28A745',      # Green
                'Awaiting QA': 'FFC107', # Yellow  
                'Awaiting Dept': 'DC3545' # Red
            }
            
            for i, series in enumerate(chart.series):
                series_name = ['Closed', 'Awaiting QA', 'Awaiting Dept'][i]
                if series_name in series_colors:
                    color_hex = series_colors[series_name]
                    r = int(color_hex[0:2], 16)
                    g = int(color_hex[2:4], 16)
                    b = int(color_hex[4:6], 16)
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = RGBColor(r, g, b)
                    print(f"  ✅ Formatted series '{series_name}' with color {color_hex}")
            
            print("✅ Successfully created new disposition chart")
            
            # Print summary
            total_missed = disposition_df_sorted[[disposition_columns['closed'], 
                                     disposition_columns['awaiting_qa'], 
                                     disposition_columns['awaiting_dept']]].sum().sum()
            print(f"📈 Disposition summary: Total missed work orders = {total_missed}")
            
        except Exception as e:
            print(f"❌ Error creating disposition chart: {e}")
            traceback.print_exc()
            
    except Exception as e:
        print(f"❌ Error updating disposition chart: {e}")
        traceback.print_exc()

def update_group_charts(prs, by_group_df, slide_index=3):
    """
    Updates group-based charts on slide 3
    """
    if by_group_df is None or by_group_df.empty:
        print("⚠️ No group data provided - skipping group charts")
        return
        
    slide = prs.slides[slide_index]
    
    # Debug: Print DataFrame info
    print("🔍 DEBUG: by_group_df columns:", by_group_df.columns.tolist())
    print("🔍 DEBUG: by_group_df shape:", by_group_df.shape)
    
    # Map column names for group charts
    group_column_mapping = {
        'group': ['group', 'Group', 'GROUP'],
        'missed': ['missed', 'Missed', 'MISSED'],
        'missed_percentage': ['missed_percentage', 'Missed %', 'missed_pct', 'Miss %'],
        'still_open': ['still_open', 'Still Open', 'Open', 'OPEN']
    }
    
    # Find actual column names
    actual_group_columns = {}
    for standard_name, possible_names in group_column_mapping.items():
        for possible_name in possible_names:
            if possible_name in by_group_df.columns:
                actual_group_columns[standard_name] = possible_name
                break
    
    print(f"📊 Group column mapping: {actual_group_columns}")
    
    # Chart titles to look for
    chart_titles = [
        "Qty Missed by Group",
        "% Missed by Group", 
        "Missed Still Open by Group"
    ]
    
    for shape in slide.shapes:
        if hasattr(shape, "chart"):
            chart = shape.chart
            chart_title = chart.chart_title.text_frame.text if chart.has_title else "Unknown Chart"
            
            if chart_title in chart_titles:
                print(f"Updating chart: {chart_title}")
                
                try:
                    # Check if we have required columns for this chart
                    if 'group' not in actual_group_columns:
                        print(f"⚠️ No group column found - skipping {chart_title}")
                        continue
                    
                    # Update chart data based on title
                    chart_data = CategoryChartData()
                    chart_data.categories = by_group_df[actual_group_columns['group']].tolist()
                    
                    if "Qty Missed" in chart_title and 'missed' in actual_group_columns:
                        chart_data.add_series('Missed', by_group_df[actual_group_columns['missed']].tolist())
                    elif "% Missed" in chart_title and 'missed_percentage' in actual_group_columns:
                        chart_data.add_series('% Missed', by_group_df[actual_group_columns['missed_percentage']].tolist())
                    elif "Still Open" in chart_title and 'still_open' in actual_group_columns:
                        chart_data.add_series('Still Open', by_group_df[actual_group_columns['still_open']].tolist())
                    else:
                        print(f"⚠️ Required columns not found for {chart_title}")
                        continue
                    
                    # Replace chart data
                    chart.replace_data(chart_data)
                    print(f"✅ Updated {chart_title}")
                    
                except Exception as e:
                    print(f"❌ Error updating {chart_title}: {e}")

def generate_summary_stats(by_month_df, disposition_df, by_group_df):
    """
    Generates summary statistics for the governance deck
    """
    try:
        # Monthly totals
        total_generated = by_month_df['generated'].sum()
        total_completed = by_month_df['completed'].sum() 
        total_missed = by_month_df['missed'].sum()
        
        # Performance metrics
        completion_rate = (total_completed / total_generated * 100) if total_generated > 0 else 0
        miss_rate = (total_missed / total_generated * 100) if total_generated > 0 else 0
        
        # Group performance
        worst_group = by_group_df.loc[by_group_df['missed_percentage'].idxmax(), 'group'] if not by_group_df.empty else 'N/A'
        best_group = by_group_df.loc[by_group_df['missed_percentage'].idxmin(), 'group'] if not by_group_df.empty else 'N/A'
        
        # Recent trend (last 3 months)
        recent_months = by_month_df.tail(3)
        recent_miss_rate = (recent_months['missed'].sum() / recent_months['generated'].sum() * 100) if recent_months['generated'].sum() > 0 else 0
        
        summary = {
            'total_generated': total_generated,
            'total_completed': total_completed,
            'total_missed': total_missed,
            'completion_rate': completion_rate,
            'miss_rate': miss_rate,
            'worst_group': worst_group,
            'best_group': best_group,
            'recent_miss_rate': recent_miss_rate
        }
        
        print("📊 Summary Statistics Generated:")
        print(f"  Total Work Orders: {total_generated:,}")
        print(f"  Completion Rate: {completion_rate:.1f}%")
        print(f"  Miss Rate: {miss_rate:.1f}%")
        print(f"  Worst Performing Group: {worst_group}")
        print(f"  Best Performing Group: {best_group}")
        print(f"  Recent 3-Month Miss Rate: {recent_miss_rate:.1f}%")
        
        return summary
        
    except Exception as e:
        print(f"❌ Error generating summary stats: {e}")
        return {}

def validate_slide_content(prs):
    """
    Validates that all slides have the expected content
    """
    try:
        print("🔍 Validating slide content...")
        
        for i, slide in enumerate(prs.slides):
            print(f"📄 Slide {i + 1}:")
            
            # Count different types of shapes
            charts = 0
            tables = 0
            text_boxes = 0
            
            for shape in slide.shapes:
                if hasattr(shape, "chart"):
                    charts += 1
                elif hasattr(shape, "table"):
                    tables += 1
                elif hasattr(shape, "text_frame"):
                    text_boxes += 1
            
            print(f"  Charts: {charts}, Tables: {tables}, Text boxes: {text_boxes}")
        
        print("✅ Slide validation complete")
        
    except Exception as e:
        print(f"❌ Error validating slides: {e}")

# Legacy function kept for backward compatibility
def create_monthly_governance_slide(by_month_df, filename=None):
    """
    Legacy function - use create_full_governance_deck instead
    """
    print("⚠️ Using legacy function - consider updating to create_full_governance_deck")
    return create_full_governance_deck(by_month_df, None, None, None, filename)

def create_governance_slide(by_month_df, filename=None):
    """
    Creates a governance slide with monthly data
    Wrapper function for backward compatibility with GUI
    """
    print("📊 Creating governance slide from GUI...")
    return create_full_governance_deck(by_month_df, None, None, None, filename)
