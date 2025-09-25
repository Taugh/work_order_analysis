# ---------------------------------------------------------------
# chart_builder.py
#
# Purpose:
#   Builds and exports charts for work order analysis, including Excel combo charts
#   and PowerPoint slide integration for preventive maintenance trends.
#
# Requirements:
#   - Input: pandas DataFrame (by_month_df) with columns: 'report_month', 'missed', 'completed', 'generated'.
#   - Libraries: pandas, xlsxwriter, pptx.
#   - Output paths: Excel file for chart data and image, PowerPoint file for slides.
#
# Output:
#   - Exports a combo chart (bar + line) to Excel in 'outputs/reports/pm_chart.xlsx'.
#   - Adds a slide with a PM chart image to a PowerPoint presentation.
#
# Notes:
#   - Used by reporting modules to visualize monthly preventive maintenance metrics.
#   - Functions: export_chart_to_excel (Excel chart), create_missed_by_month_slide (PowerPoint slide).
# ---------------------------------------------------------------

# scripts/chart_builder.py

from pptx.util import Inches
import pandas as pd
import os

def export_chart_to_excel(by_month_df, filename="outputs/reports/pm_chart.xlsx"):
    """
    Exports a combo chart (bar + line) to an Excel file using by_month_df.
    - Missed: bar chart
    - Completed: line chart
    - Generated: line chart
    """
    os.makedirs(os.path.dirname(filename), exist_ok=True)

    with pd.ExcelWriter(filename, engine="xlsxwriter") as writer:
        by_month_df.to_excel(writer, index=False, sheet_name="ChartData")
        workbook = writer.book
        worksheet = writer.sheets["ChartData"]

        chart = workbook.add_chart({'type': 'column'})

        # Bar series: Missed
        chart.add_series({
            'name':       'Missed',
            'categories': ['ChartData', 1, 0, len(by_month_df), 0],  # report_month
            'values':     ['ChartData', 1, 1, len(by_month_df), 1],  # missed
        })

        # Line series: Completed
        chart.add_series({
            'name':       'Completed',
            'categories': ['ChartData', 1, 0, len(by_month_df), 0],
            'values':     ['ChartData', 1, 2, len(by_month_df), 2],
            'type':       'line',
            'y2_axis':    True,
            'marker':     {'type': 'circle', 'size': 5},
        })

        # Line series: Generated
        chart.add_series({
            'name':       'Generated',
            'categories': ['ChartData', 1, 0, len(by_month_df), 0],
            'values':     ['ChartData', 1, 3, len(by_month_df), 3],
            'type':       'line',
            'y2_axis':    True,
            'marker':     {'type': 'square', 'size': 5},
        })

        # Chart formatting
        chart.set_title({'name': 'Preventive Maintenance Trends'})
        chart.set_x_axis({'name': 'Month'})
        chart.set_y_axis({'name': 'Missed'})
        chart.set_y2_axis({'name': 'Completed / Generated'})
        chart.set_legend({'position': 'top'})

        worksheet.insert_chart('F2', chart)

    print(f"📊 Excel chart exported to: {filename}")

def create_missed_by_month_slide(prs):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Preventive Maintenance Trends"

    img_path = "outputs/reports/pm_chart.png"
    slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))
